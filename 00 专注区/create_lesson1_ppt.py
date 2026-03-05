from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# 定义颜色
COLOR_PRIMARY = RGBColor(28, 40, 51)
COLOR_ACCENT = RGBColor(52, 152, 219)
COLOR_QUESTION = RGBColor(241, 196, 15)
COLOR_EXERCISE = RGBColor(232, 248, 245)
COLOR_ANSWER = RGBColor(213, 244, 230)
COLOR_RED = RGBColor(231, 76, 60)
COLOR_GREEN = RGBColor(39, 174, 96)

def add_title_slide(prs, title, subtitle=""):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 设置背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

def add_question_slide(prs, question):
    """添加提问页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 设置背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_QUESTION
    
    # 添加问号
    question_mark = slide.shapes.add_textbox(Inches(4), Inches(1), Inches(2), Inches(2))
    frame = question_mark.text_frame
    frame.text = "?"
    p = frame.paragraphs[0]
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 添加问题文本
    question_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
    frame = question_box.text_frame
    frame.text = question
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(44, 62, 80)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    frame = title_box.text_frame
    frame.text = title
    p = frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # 添加内容
    y_pos = 1.2
    for item in content_list:
        content_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.5))
        frame = content_box.text_frame
        frame.text = item
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.font.size = Pt(18)
        y_pos += 0.6

def add_exercise_slide(prs, title, questions, time="2 分钟"):
    """添加练习题目页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 设置背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_EXERCISE
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(0.6))
    frame = title_box.text_frame
    frame.text = title
    p = frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # 添加时间
    time_box = slide.shapes.add_textbox(Inches(7), Inches(0.5), Inches(2.5), Inches(0.5))
    frame = time_box.text_frame
    frame.text = f"⏱ 答题时间：{time}"
    p = frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_RED
    p.alignment = PP_ALIGN.RIGHT
    
    # 添加题目
    y_pos = 1.8
    for q in questions:
        q_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.6))
        frame = q_box.text_frame
        frame.text = q
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.font.size = Pt(22)
        y_pos += 1

def add_answer_slide(prs, title, answers):
    """添加答案页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 设置背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_ANSWER
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    frame = title_box.text_frame
    frame.text = title
    p = frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # 添加答案
    y_pos = 1.5
    for ans in answers:
        ans_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.7))
        frame = ans_box.text_frame
        frame.text = ans['text']
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        if ans.get('correct'):
            p.font.color.rgb = COLOR_GREEN
        else:
            p.font.color.rgb = COLOR_RED
        
        # 添加注释
        if ans.get('note'):
            note_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.4), Inches(8), Inches(0.3))
            frame = note_box.text_frame
            frame.text = f"   {ans['note']}"
            p = frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(127, 140, 141)
            p.font.italic = True
        
        y_pos += 1

# ========== 开始生成 PPT ==========

# 第 1 页：封面
add_title_slide(prs, "数字电路基础 + 与门", "第 1 课")

# 第 2 页：提问 1
add_question_slide(prs, '生活中哪些情况是\n"同时满足"才能发生的？\n\n请举 2-3 个例子')

# 第 3 页：数字电路基础
add_content_slide(prs, "一、数字电路基础", [
    "电子线路中的电信号有两大类：模拟信号 和 数字信号",
    "",
    "• 模拟信号：在数值上和时间上都是连续变化的信号",
    "• 数字信号：在数值上和时间上不连续变化的信号",
    "• 模拟电路：处理模拟信号的电路",
    "• 数字电路：处理数字信号的电路"
])

# 第 4 页：逻辑状态
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
frame = title_box.text_frame
frame.text = "二、逻辑状态的表示"
p = frame.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.8))
frame = subtitle_box.text_frame
frame.text = "用数字符号 0 和 1 表示相互对立的逻辑状态"
p = frame.paragraphs[0]
p.font.size = Pt(24)
p.font.color.rgb = COLOR_ACCENT
p.alignment = PP_ALIGN.CENTER

# 添加逻辑 0 和 1 的框
left_box = slide.shapes.add_shape(1, Inches(2), Inches(2.5), Inches(2.5), Inches(1.5))
left_box.fill.solid()
left_box.fill.fore_color.rgb = RGBColor(232, 248, 245)
left_text = left_box.text_frame
left_text.text = "逻辑 0"
p = left_text.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

right_box = slide.shapes.add_shape(1, Inches(5.5), Inches(2.5), Inches(2.5), Inches(1.5))
right_box.fill.solid()
right_box.fill.fore_color.rgb = RGBColor(232, 248, 245)
right_text = right_box.text_frame
right_text.text = "逻辑 1"
p = right_text.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# 第 5 页：提问 2
add_question_slide(prs, "如果用 0 和 1 表示开关的状态\n开关打开用什么表示？\n关闭呢？")

# 第 6 页：高低电平
add_content_slide(prs, "三、高低电平", [
    "用高电平、低电平来描述电位的高低",
    "",
    "重要：高低电平不是固定值，而是电平变化范围",
    "",
    "• 标准高电平 VSH —— 高电平的下限值",
    "• 标准低电平 VSL —— 低电平的上限值",
    "• 应用时：高电平 ≥ VSH，低电平 ≤ VSL"
])

# 第 7 页：练习 0.1（题目）
add_exercise_slide(prs, "练习 0.1：判断题", [
    "1. 高电平一定是 5V。",
    "2. 低电平一定小于或等于 VSL。",
    "3. 高电平和低电平是固定值。"
])

# 第 8 页：练习 0.1（答案）
add_answer_slide(prs, "练习 0.1：答案", [
    {'text': '1. 高电平一定是 5V。 ×', 'correct': False, 'note': '高电平是范围，不是固定值'},
    {'text': '2. 低电平一定小于或等于 VSL。 √', 'correct': True, 'note': '正确'},
    {'text': '3. 高电平和低电平是固定值。 ×', 'correct': False, 'note': '是范围，不是固定值'}
])

# 第 9 页：正负逻辑
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
frame = title_box.text_frame
frame.text = "四、正逻辑和负逻辑"
p = frame.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

# 正逻辑框
pos_box = slide.shapes.add_shape(1, Inches(1), Inches(1.5), Inches(3.5), Inches(2))
pos_box.fill.solid()
pos_box.fill.fore_color.rgb = RGBColor(232, 248, 245)
pos_text = pos_box.text_frame
pos_text.text = "正逻辑\n\n1 = 高电平\n0 = 低电平"
for p in pos_text.paragraphs:
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER
pos_text.paragraphs[0].font.size = Pt(24)
pos_text.paragraphs[0].font.bold = True
pos_text.paragraphs[0].font.color.rgb = COLOR_ACCENT

# 负逻辑框
neg_box = slide.shapes.add_shape(1, Inches(5.5), Inches(1.5), Inches(3.5), Inches(2))
neg_box.fill.solid()
neg_box.fill.fore_color.rgb = RGBColor(252, 243, 207)
neg_text = neg_box.text_frame
neg_text.text = "负逻辑\n\n1 = 低电平\n0 = 高电平"
for p in neg_text.paragraphs:
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER
neg_text.paragraphs[0].font.size = Pt(24)
neg_text.paragraphs[0].font.bold = True
neg_text.paragraphs[0].font.color.rgb = RGBColor(243, 156, 18)

# 添加说明
note_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
frame = note_box.text_frame
frame.text = "本课程默认使用正逻辑"
p = frame.paragraphs[0]
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLOR_RED
p.alignment = PP_ALIGN.CENTER

# 第 10 页：练习 0.2（题目）
add_exercise_slide(prs, "练习 0.2：填空题", [
    "1. 正逻辑规定：用 ___ 表示高电平，用 ___ 表示低电平。",
    "2. 负逻辑规定：用 ___ 表示低电平，用 ___ 表示高电平。"
])

# 第 11 页：练习 0.2（答案）
add_answer_slide(prs, "练习 0.2：答案", [
    {'text': '1. 正逻辑规定：用 1 表示高电平，用 0 表示低电平。', 'correct': True},
    {'text': '2. 负逻辑规定：用 1 表示低电平，用 0 表示高电平。', 'correct': True}
])

# 第 12 页：提问 3
add_question_slide(prs, '"两个开关都打开，灯才亮"\n\n这是什么逻辑关系？')

# 第 13 页：与门定义
add_content_slide(prs, "五、与门（AND Gate）", [
    "1. 与逻辑关系",
    "   当决定一件事情的几个条件全部具备后，这件事情才能发生。",
    "",
    "2. 逻辑符号：&",
    "",
    "3. 逻辑函数式：Y = A · B  或  Y = AB"
])

# 第 14 页：与门真值表
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
frame = title_box.text_frame
frame.text = "与门真值表"
p = frame.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

# 添加真值表
table_data = [
    ["A", "B", "Y"],
    ["0", "0", "0"],
    ["0", "1", "0"],
    ["1", "0", "0"],
    ["1", "1", "1"]
]

table = slide.shapes.add_table(5, 3, Inches(2.5), Inches(1.5), Inches(5), Inches(2.5)).table

for i, row_data in enumerate(table_data):
    for j, cell_data in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_data
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(24)
            paragraph.alignment = PP_ALIGN.CENTER
            if i == 0:  # 表头
                paragraph.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_ACCENT
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
            elif i == 4 and j == 2:  # 最后一行的输出
                paragraph.font.bold = True
                paragraph.font.color.rgb = COLOR_GREEN

# 添加口诀
formula_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.5))
frame = formula_box.text_frame
frame.text = "逻辑函数式：Y = A · B  或  Y = AB"
p = frame.paragraphs[0]
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT
p.alignment = PP_ALIGN.CENTER

slogan_box = slide.shapes.add_textbox(Inches(1), Inches(4.7), Inches(8), Inches(0.4))
frame = slogan_box.text_frame
frame.text = "口诀：全 1 出 1，有 0 出 0"
p = frame.paragraphs[0]
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_RED
p.alignment = PP_ALIGN.CENTER

# 第 15 页：提问 4
add_question_slide(prs, "如果 A=1, B=0\n\n与门的输出 Y 是多少？")

print("正在生成 PPT，已完成前 15 页...")

# 保存文件
prs.save("第1课-数字电路基础与与门.pptx")
print("PPT 生成成功！文件：第1课-数字电路基础与与门.pptx")
print("共生成 15 页（封面 + 提问页 + 内容页 + 练习页）")
